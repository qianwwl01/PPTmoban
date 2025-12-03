# -*- coding: utf-8 -*-
"""
PPT生成器模块
封装所有与python-pptx相关的PPT生成逻辑
"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config_presets import SLIDE_RATIOS


def hex_to_rgb(hex_color: str) -> RGBColor:
    """
    将十六进制颜色转换为RGBColor对象
    
    参数:
        hex_color: 十六进制颜色字符串，如 "#1a365d"
    返回:
        RGBColor对象
    """
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def set_shape_fill(shape, color_hex: str):
    """
    设置形状的填充颜色
    
    参数:
        shape: pptx形状对象
        color_hex: 十六进制颜色
    """
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def set_text_style(text_frame, text: str, font_name: str, font_size: int, 
                   color_hex: str, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT):
    """
    设置文本框的文字样式
    
    参数:
        text_frame: 文本框对象
        text: 文本内容
        font_name: 字体名称
        font_size: 字体大小（磅）
        color_hex: 字体颜色
        bold: 是否加粗
        align: 对齐方式
    """
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(color_hex)
    p.font.bold = bold
    p.alignment = align


def add_text_box(slide, left: float, top: float, width: float, height: float,
                 text: str, font_name: str, font_size: int, color_hex: str,
                 bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT,
                 vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP):
    """
    在幻灯片上添加文本框
    
    参数:
        slide: 幻灯片对象
        left, top, width, height: 位置和尺寸（英寸）
        text: 文本内容
        font_name: 字体名称
        font_size: 字体大小
        color_hex: 字体颜色
        bold: 是否加粗
        align: 水平对齐
        vertical_anchor: 垂直对齐
    返回:
        创建的文本框形状
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    # 设置垂直对齐
    tf.anchor = vertical_anchor
    
    set_text_style(tf, text, font_name, font_size, color_hex, bold, align)
    return txBox


def add_rectangle(slide, left: float, top: float, width: float, height: float, 
                  fill_color: str, line_color: str = None):
    """
    在幻灯片上添加矩形
    
    参数:
        slide: 幻灯片对象
        left, top, width, height: 位置和尺寸（英寸）
        fill_color: 填充颜色
        line_color: 边框颜色（可选）
    返回:
        创建的矩形形状
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_shape_fill(shape, fill_color)
    
    if line_color:
        shape.line.color.rgb = hex_to_rgb(line_color)
    else:
        shape.line.fill.background()
    
    return shape


def set_slide_background(slide, color_hex: str):
    """
    设置幻灯片背景颜色
    
    参数:
        slide: 幻灯片对象
        color_hex: 背景颜色
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_title_slide(prs: Presentation, config: dict):
    """
    添加标题页
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    set_slide_background(slide, config['background'])
    
    # 获取幻灯片尺寸
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 顶部装饰条
    add_rectangle(slide, 0, 0, slide_width, 0.15, config['primary'])
    
    # 主标题
    add_text_box(
        slide, 0.5, slide_height * 0.35, slide_width - 1, 1.2,
        "在此输入演示文稿标题",
        config['title_font'], 44, config['primary'],
        bold=True, align=PP_ALIGN.CENTER
    )
    
    # 副标题
    add_text_box(
        slide, 0.5, slide_height * 0.55, slide_width - 1, 0.8,
        "在此输入副标题或简短描述",
        config['body_font'], 24, config['secondary'],
        align=PP_ALIGN.CENTER
    )
    
    # 底部信息栏
    add_rectangle(slide, 0, slide_height - 0.8, slide_width, 0.8, config['primary'])
    add_text_box(
        slide, 0.5, slide_height - 0.6, slide_width - 1, 0.4,
        "演讲者姓名  |  公司名称  |  日期",
        config['body_font'], 14, "#ffffff",
        align=PP_ALIGN.CENTER
    )


def add_agenda_slide(prs: Presentation, config: dict):
    """
    添加目录页
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 左侧装饰条
    add_rectangle(slide, 0, 0, 0.15, slide_height, config['primary'])
    
    # 页面标题
    add_text_box(
        slide, 0.8, 0.5, slide_width - 1.5, 0.8,
        "目 录",
        config['title_font'], 36, config['primary'],
        bold=True
    )
    
    # 分隔线
    add_rectangle(slide, 0.8, 1.3, 2, 0.05, config['accent'])
    
    # 目录条目
    agenda_items = [
        "01  第一部分标题",
        "02  第二部分标题", 
        "03  第三部分标题",
        "04  第四部分标题",
        "05  第五部分标题"
    ]
    
    start_y = 1.8
    for i, item in enumerate(agenda_items):
        # 条目背景
        if i % 2 == 0:
            add_rectangle(slide, 0.8, start_y + i * 0.9, slide_width - 1.6, 0.8, "#f8f9fa")
        
        add_text_box(
            slide, 1.0, start_y + i * 0.9 + 0.2, slide_width - 2, 0.5,
            item,
            config['body_font'], 20, config['secondary']
        )


def add_content_slide(prs: Presentation, config: dict, page_num: int = 1):
    """
    添加内容页
    
    参数:
        prs: Presentation对象
        config: 配置字典
        page_num: 页码（用于区分不同内容页）
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 顶部标题区
    add_rectangle(slide, 0, 0, slide_width, 1.2, config['primary'])
    add_text_box(
        slide, 0.5, 0.35, slide_width - 1, 0.6,
        f"内容页标题 - 第{page_num}页",
        config['title_font'], 32, "#ffffff",
        bold=True
    )
    
    # 内容区域
    content_text = """• 在此输入第一个要点内容
    
• 在此输入第二个要点内容
    - 子要点说明文字
    - 更多细节描述
    
• 在此输入第三个要点内容

• 在此输入第四个要点内容"""
    
    add_text_box(
        slide, 0.8, 1.6, slide_width - 1.6, slide_height - 2.5,
        content_text,
        config['body_font'], 18, config['secondary']
    )
    
    # 底部页码
    add_text_box(
        slide, slide_width - 1.5, slide_height - 0.5, 1, 0.3,
        f"第 {len(prs.slides)} 页",
        config['body_font'], 10, config['secondary'],
        align=PP_ALIGN.RIGHT
    )


def add_image_text_slide(prs: Presentation, config: dict, layout_variant: str = 'left-image', image_bytes: bytes = None):
    """
    添加图文页
    
    参数:
        prs: Presentation对象
        config: 配置字典
        layout_variant: 布局变体 ('left-image' 或 'right-image')
        image_bytes: 图片字节数据（可选）
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 页面标题
    add_text_box(
        slide, 0.5, 0.3, slide_width - 1, 0.7,
        "图文混排页标题",
        config['title_font'], 28, config['primary'],
        bold=True
    )
    
    # 标题下划线
    add_rectangle(slide, 0.5, 1.0, 3, 0.05, config['accent'])
    
    content_y = 1.3
    content_height = slide_height - 1.8
    
    if layout_variant == 'left-image':
        # 左图右文布局
        img_left = 0.5
        img_width = 5.5
        
        # 如果有图片，插入真实图片
        if image_bytes:
            try:
                img_stream = io.BytesIO(image_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(img_left), Inches(content_y),
                    width=Inches(img_width)
                )
            except Exception:
                # 图片插入失败，显示占位区
                add_rectangle(slide, img_left, content_y, img_width, content_height, "#e2e8f0", config['secondary'])
                add_text_box(slide, img_left, content_y + content_height/2 - 0.3, img_width, 0.6,
                    "📷 图片占位区域", config['body_font'], 16, config['secondary'], align=PP_ALIGN.CENTER)
        else:
            # 没有图片，显示占位区
            add_rectangle(slide, img_left, content_y, img_width, content_height, "#e2e8f0", config['secondary'])
            add_text_box(slide, img_left, content_y + content_height/2 - 0.3, img_width, 0.6,
                "📷 图片占位区域\n点击添加图片", config['body_font'], 16, config['secondary'], align=PP_ALIGN.CENTER)
        
        # 右侧文字
        text_content = """在此输入说明文字

• 要点一：详细描述内容

• 要点二：详细描述内容

• 要点三：详细描述内容

可以在这里添加更多的解释性文字来配合左侧的图片内容。"""
        
        add_text_box(
            slide, 6.3, content_y + 0.2, slide_width - 7, content_height - 0.4,
            text_content,
            config['body_font'], 16, config['secondary']
        )
    else:
        # 右图左文布局
        # 左侧文字
        text_content = """在此输入说明文字

• 要点一：详细描述内容

• 要点二：详细描述内容

• 要点三：详细描述内容

可以在这里添加更多的解释性文字来配合右侧的图片内容。"""
        
        add_text_box(
            slide, 0.5, content_y + 0.2, 5.5, content_height - 0.4,
            text_content,
            config['body_font'], 16, config['secondary']
        )
        
        # 右侧图片区
        img_left = 6.3
        img_width = slide_width - 6.8
        
        if image_bytes:
            try:
                img_stream = io.BytesIO(image_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(img_left), Inches(content_y),
                    width=Inches(img_width)
                )
            except Exception:
                add_rectangle(slide, img_left, content_y, img_width, content_height, "#e2e8f0", config['secondary'])
                add_text_box(slide, img_left, content_y + content_height/2 - 0.3, img_width, 0.6,
                    "📷 图片占位区域", config['body_font'], 16, config['secondary'], align=PP_ALIGN.CENTER)
        else:
            add_rectangle(slide, img_left, content_y, img_width, content_height, "#e2e8f0", config['secondary'])
            add_text_box(slide, img_left, content_y + content_height/2 - 0.3, img_width, 0.6,
                "📷 图片占位区域\n点击添加图片", config['body_font'], 16, config['secondary'], align=PP_ALIGN.CENTER)


def add_comparison_slide(prs: Presentation, config: dict):
    """
    添加对比页
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 页面标题
    add_text_box(
        slide, 0.5, 0.3, slide_width - 1, 0.7,
        "对比分析页",
        config['title_font'], 28, config['primary'],
        bold=True, align=PP_ALIGN.CENTER
    )
    
    # 中间分隔线
    add_rectangle(slide, slide_width/2 - 0.02, 1.2, 0.04, slide_height - 1.7, config['accent'])
    
    # 左侧区块
    left_width = slide_width/2 - 0.8
    
    add_rectangle(slide, 0.4, 1.3, left_width, 0.6, config['primary'])
    add_text_box(
        slide, 0.4, 1.4, left_width, 0.4,
        "方案 A",
        config['title_font'], 20, "#ffffff",
        bold=True, align=PP_ALIGN.CENTER
    )
    
    left_content = """✓ 优势点一

✓ 优势点二

✓ 优势点三

✗ 不足之处"""
    
    add_text_box(
        slide, 0.5, 2.1, left_width - 0.2, slide_height - 2.8,
        left_content,
        config['body_font'], 16, config['secondary']
    )
    
    # 右侧区块
    right_x = slide_width/2 + 0.3
    
    add_rectangle(slide, right_x, 1.3, left_width, 0.6, config['accent'])
    add_text_box(
        slide, right_x, 1.4, left_width, 0.4,
        "方案 B",
        config['title_font'], 20, "#ffffff",
        bold=True, align=PP_ALIGN.CENTER
    )
    
    right_content = """✓ 优势点一

✓ 优势点二

✓ 优势点三

✗ 不足之处"""
    
    add_text_box(
        slide, right_x + 0.1, 2.1, left_width - 0.2, slide_height - 2.8,
        right_content,
        config['body_font'], 16, config['secondary']
    )


def add_thankyou_slide(prs: Presentation, config: dict):
    """
    添加致谢页
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 装饰圆形（或矩形模拟）
    center_x = slide_width / 2
    center_y = slide_height / 2
    
    # 背景装饰
    add_rectangle(slide, center_x - 4, center_y - 1.5, 8, 3, config['primary'])
    
    # 主标题
    add_text_box(
        slide, 0.5, center_y - 0.8, slide_width - 1, 1,
        "感谢观看",
        config['title_font'], 48, "#ffffff",
        bold=True, align=PP_ALIGN.CENTER
    )
    
    # 副文本
    add_text_box(
        slide, 0.5, center_y + 0.3, slide_width - 1, 0.6,
        "THANK YOU FOR WATCHING",
        config['body_font'], 18, "#ffffff",
        align=PP_ALIGN.CENTER
    )
    
    # 底部联系信息
    add_text_box(
        slide, 0.5, slide_height - 1, slide_width - 1, 0.5,
        "联系方式：email@example.com  |  电话：123-4567-8900",
        config['body_font'], 12, config['secondary'],
        align=PP_ALIGN.CENTER
    )


def add_timeline_slide(prs: Presentation, config: dict):
    """
    添加时间轴页
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 页面标题
    add_text_box(
        slide, 0.5, 0.3, slide_width - 1, 0.7,
        "项目时间轴 / 里程碑",
        config['title_font'], 28, config['primary'],
        bold=True
    )
    
    # 标题下划线
    add_rectangle(slide, 0.5, 1.0, 3, 0.05, config['accent'])
    
    # 时间轴主线
    timeline_y = slide_height / 2
    add_rectangle(slide, 0.8, timeline_y - 0.03, slide_width - 1.6, 0.06, config['primary'])
    
    # 时间节点
    nodes = [
        ("2024 Q1", "第一阶段\n项目启动"),
        ("2024 Q2", "第二阶段\n设计开发"),
        ("2024 Q3", "第三阶段\n测试优化"),
        ("2024 Q4", "第四阶段\n正式上线")
    ]
    
    node_spacing = (slide_width - 2) / (len(nodes) + 1)
    
    for i, (date, desc) in enumerate(nodes):
        x = 1 + node_spacing * (i + 1) - 0.4
        
        # 节点圆圈
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(timeline_y - 0.2), Inches(0.4), Inches(0.4)
        )
        set_shape_fill(shape, config['accent'])
        shape.line.fill.background()
        
        # 日期标签（上方）
        add_text_box(
            slide, x - 0.3, timeline_y - 0.9, 1, 0.5,
            date,
            config['body_font'], 14, config['primary'],
            bold=True, align=PP_ALIGN.CENTER
        )
        
        # 描述文字（下方）
        add_text_box(
            slide, x - 0.5, timeline_y + 0.4, 1.4, 0.8,
            desc,
            config['body_font'], 12, config['secondary'],
            align=PP_ALIGN.CENTER
        )


def add_kpi_slide(prs: Presentation, config: dict):
    """
    添加数据概览页 (KPI展示)
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 页面标题
    add_text_box(
        slide, 0.5, 0.3, slide_width - 1, 0.7,
        "核心数据概览",
        config['title_font'], 28, config['primary'],
        bold=True, align=PP_ALIGN.CENTER
    )
    
    # KPI 卡片
    kpis = [
        ("1,234", "总用户数", "+12.5%"),
        ("98.6%", "系统可用率", "+2.1%"),
        ("56.7万", "月访问量", "+25.3%"),
        ("4.8/5", "用户满意度", "+0.3")
    ]
    
    card_width = (slide_width - 1.5) / 4
    card_height = 2.5
    start_y = (slide_height - card_height) / 2
    
    for i, (number, label, change) in enumerate(kpis):
        x = 0.5 + i * (card_width + 0.15)
        
        # 卡片背景
        add_rectangle(slide, x, start_y, card_width - 0.1, card_height, "#f8f9fa", config['secondary'])
        
        # 数字
        add_text_box(
            slide, x, start_y + 0.4, card_width - 0.1, 0.8,
            number,
            config['title_font'], 36, config['primary'],
            bold=True, align=PP_ALIGN.CENTER
        )
        
        # 标签
        add_text_box(
            slide, x, start_y + 1.2, card_width - 0.1, 0.5,
            label,
            config['body_font'], 14, config['secondary'],
            align=PP_ALIGN.CENTER
        )
        
        # 增长标记
        change_color = "#10b981" if change.startswith("+") else "#ef4444"
        add_text_box(
            slide, x, start_y + 1.7, card_width - 0.1, 0.4,
            change,
            config['body_font'], 14, change_color,
            bold=True, align=PP_ALIGN.CENTER
        )


def add_quote_slide(prs: Presentation, config: dict):
    """
    添加引用页
    
    参数:
        prs: Presentation对象
        config: 配置字典
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_slide_background(slide, config['background'])
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    center_y = slide_height / 2
    
    # 左侧装饰竖线
    add_rectangle(slide, 1, center_y - 1.5, 0.1, 3, config['accent'])
    
    # 引号装饰
    add_text_box(
        slide, 1.3, center_y - 1.8, 1, 1,
        "“",
        config['title_font'], 72, config['accent'],
        bold=True
    )
    
    # 引用文字
    add_text_box(
        slide, 1.5, center_y - 0.8, slide_width - 3, 1.6,
        "在此输入引言或重要语句，\n用于强调核心观点或名人名言。",
        config['body_font'], 28, config['primary'],
        align=PP_ALIGN.LEFT
    )
    
    # 作者/来源
    add_text_box(
        slide, 1.5, center_y + 1.2, slide_width - 3, 0.5,
        "—— 作者姓名，《来源出处》",
        config['body_font'], 16, config['secondary'],
        align=PP_ALIGN.LEFT
    )


def add_watermark(slide, text: str, opacity: int, slide_width: float, slide_height: float):
    """
    在幻灯片中央添加水印
    
    参数:
        slide: 幻灯片对象
        text: 水印文字
        opacity: 透明度 (0-100)
        slide_width, slide_height: 幻灯片尺寸
    """
    # 计算透明度对应的颜色值
    gray_value = 255 - int(opacity * 2.55)
    color_hex = f"#{gray_value:02x}{gray_value:02x}{gray_value:02x}"
    
    txBox = slide.shapes.add_textbox(
        Inches(0), Inches(slide_height / 2 - 0.5),
        Inches(slide_width), Inches(1)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(48)
    p.font.color.rgb = hex_to_rgb(color_hex)
    p.alignment = PP_ALIGN.CENTER


def add_logo_to_slide(slide, logo_bytes: bytes, slide_width: float, slide_height: float, position: str = "bottom-right"):
    """
    在幻灯片上添加Logo
    
    参数:
        slide: 幻灯片对象
        logo_bytes: Logo图片的字节数据
        slide_width, slide_height: 幻灯片尺寸
        position: 位置 ('bottom-right', 'bottom-left', 'top-right', 'top-left')
    """
    logo_stream = io.BytesIO(logo_bytes)
    logo_height = 0.5  # Logo高度（英寸）
    
    # 根据位置计算坐标
    positions = {
        "bottom-right": (slide_width - 1.5, slide_height - 0.7),
        "bottom-left": (0.3, slide_height - 0.7),
        "top-right": (slide_width - 1.5, 0.2),
        "top-left": (0.3, 0.2)
    }
    
    left, top = positions.get(position, positions["bottom-right"])
    
    slide.shapes.add_picture(
        logo_stream,
        Inches(left), Inches(top),
        height=Inches(logo_height)
    )


def add_footer(slide, config: dict, slide_num: int, slide_width: float, slide_height: float):
    """
    添加页脚（页码和页脚文字）
    
    参数:
        slide: 幻灯片对象
        config: 配置字典
        slide_num: 当前页码
        slide_width, slide_height: 幻灯片尺寸
    """
    footer_y = slide_height - 0.4
    
    # 页脚文字（左侧）
    footer_text = config.get('footer_text', '')
    if footer_text:
        add_text_box(
            slide, 0.3, footer_y, 4, 0.3,
            footer_text,
            config['body_font'], 9, config['secondary'],
            align=PP_ALIGN.LEFT
        )
    
    # 页码（右侧）
    if config.get('show_page_number', True):
        add_text_box(
            slide, slide_width - 1, footer_y, 0.7, 0.3,
            str(slide_num),
            config['body_font'], 10, config['secondary'],
            align=PP_ALIGN.RIGHT
        )


def build_presentation(config: dict, layouts_config: dict, logo_bytes: bytes = None, uploaded_images: list = None) -> io.BytesIO:
    """
    根据配置生成完整的PPT模板
    
    参数:
        config: 主题配置字典，包含颜色、字体等
        layouts_config: 版式配置，指定每种版式的启用状态和数量
        logo_bytes: Logo图片字节数据（可选）
        uploaded_images: 上传的图片列表（可选）
    
    返回:
        包含PPT文件的BytesIO对象
    """
    if uploaded_images is None:
        uploaded_images = []
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸
    ratio = config.get('ratio', '16:9')
    ratio_config = SLIDE_RATIOS.get(ratio, SLIDE_RATIOS['16:9'])
    prs.slide_width = Inches(ratio_config['width'])
    prs.slide_height = Inches(ratio_config['height'])
    
    slide_width = ratio_config['width']
    slide_height = ratio_config['height']
    
    # 根据配置添加各类幻灯片
    
    # 标题页
    if layouts_config.get('title', {}).get('enabled', True):
        count = layouts_config['title'].get('count', 1)
        for _ in range(count):
            add_title_slide(prs, config)
    
    # 目录页
    if layouts_config.get('agenda', {}).get('enabled', True):
        count = layouts_config['agenda'].get('count', 1)
        for _ in range(count):
            add_agenda_slide(prs, config)
    
    # 内容页
    if layouts_config.get('content', {}).get('enabled', True):
        count = layouts_config['content'].get('count', 2)
        for i in range(count):
            add_content_slide(prs, config, page_num=i+1)
    
    # 图文页
    if layouts_config.get('image_text', {}).get('enabled', True):
        count = layouts_config['image_text'].get('count', 2)
        for i in range(count):
            variant = 'left-image' if i % 2 == 0 else 'right-image'
            # 获取对应的图片
            image_bytes = None
            if i < len(uploaded_images):
                image_bytes = uploaded_images[i].get('bytes')
            add_image_text_slide(prs, config, layout_variant=variant, image_bytes=image_bytes)
    
    # 对比页
    if layouts_config.get('comparison', {}).get('enabled', True):
        count = layouts_config['comparison'].get('count', 1)
        for _ in range(count):
            add_comparison_slide(prs, config)
    
    # 时间轴页
    if layouts_config.get('timeline', {}).get('enabled', True):
        count = layouts_config['timeline'].get('count', 1)
        for _ in range(count):
            add_timeline_slide(prs, config)
    
    # 数据概览页
    if layouts_config.get('kpi', {}).get('enabled', True):
        count = layouts_config['kpi'].get('count', 1)
        for _ in range(count):
            add_kpi_slide(prs, config)
    
    # 引用页
    if layouts_config.get('quote', {}).get('enabled', True):
        count = layouts_config['quote'].get('count', 1)
        for _ in range(count):
            add_quote_slide(prs, config)
    
    # 致谢页
    if layouts_config.get('thankyou', {}).get('enabled', True):
        count = layouts_config['thankyou'].get('count', 1)
        for _ in range(count):
            add_thankyou_slide(prs, config)
    
    # 为所有幻灯片添加水印、Logo、页脚
    for idx, slide in enumerate(prs.slides):
        # 添加水印
        if config.get('watermark_enabled', False):
            watermark_text = config.get('watermark_text', '内部资料')
            watermark_opacity = config.get('watermark_opacity', 15)
            add_watermark(slide, watermark_text, watermark_opacity, slide_width, slide_height)
        
        # 添加Logo
        if logo_bytes:
            try:
                add_logo_to_slide(slide, logo_bytes, slide_width, slide_height, "bottom-right")
            except Exception:
                pass  # 如果Logo添加失败，静默跳过
        
        # 添加页脚（跳过第一页标题页）
        if idx > 0:
            add_footer(slide, config, idx + 1, slide_width, slide_height)
    
    # 保存到内存
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return ppt_buffer
